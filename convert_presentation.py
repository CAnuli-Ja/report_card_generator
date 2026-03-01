#!/usr/bin/env python3
"""
Convert PRESENTATION.html to PowerPoint (.pptx) format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_powerpoint():
    """Create a PowerPoint presentation from slide content"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color schemes for each slide
    slide_colors = [
        ((102, 126, 234), (255, 255, 255)),  # Slide 1: Purple
        ((240, 147, 251), (255, 255, 255)),  # Slide 2: Pink
        ((79, 172, 254), (255, 255, 255)),   # Slide 3: Blue
        ((67, 234, 123), (255, 255, 255)),   # Slide 4: Green
        ((250, 154, 158), (51, 51, 51)),     # Slide 5: Coral
        ((48, 207, 208), (255, 255, 255)),   # Slide 6: Teal
        ((168, 237, 234), (51, 51, 51)),     # Slide 7: Light Teal
        ((255, 154, 86), (255, 255, 255)),   # Slide 8: Orange
    ]
    
    # Slide 1: Title
    add_title_slide(prs, "📋 Report Card Generator", 
                   "Automated Report Card Creation Made Simple",
                   "Transform hours of manual work into minutes",
                   slide_colors[0])
    
    # Slide 2: The Problem
    add_content_slide(prs, "The Challenge",
                     ["Creating individual report cards is time-consuming",
                      "Manual copying and pasting from Excel to Word",
                      "High risk of errors and inconsistencies",
                      "Repetitive work that takes valuable time away from educators"],
                     slide_colors[1])
    
    # Slide 3: The Solution
    add_solution_slide(prs, "✨ The Solution",
                      "One click to generate 100+ professional report cards automatically",
                      ["Upload Excel", "Select Sheet", "Upload Template", "Generate!"],
                      slide_colors[2])
    
    # Slide 4: Key Features
    add_content_slide(prs, "🎯 Key Features",
                     ["Works with your existing Excel spreadsheets",
                      "Customizable Word templates (35+ fields available)",
                      "Batch processing - generate all reports at once",
                      "Download individually or as ZIP file",
                      "No installation needed - use from your browser"],
                     slide_colors[3])
    
    # Slide 5: How to Use - Steps 1-2
    add_two_column_slide(prs, "📝 How to Use (Step 1-2)",
                        "Step 1: Upload Excel File",
                        "Click 'Choose Excel File' and select your spreadsheet with student data\n\nNote: First column should have student numbers",
                        "Step 2: Select Sheet",
                        "Choose the class/form sheet from your Excel file\n\nExample: F1J, F2W, F3J, F4W",
                        slide_colors[4])
    
    # Slide 6: How to Use - Steps 3-4
    add_two_column_slide(prs, "📝 How to Use (Step 3-4)",
                        "Step 3: Upload Template",
                        "Choose your Word report card template (.docx file)\n\nTip: Use placeholders like {{NAME}}, {{MATH}}, {{COMMENTS}}",
                        "Step 4: Generate & Download",
                        "Click the Generate button and wait for processing\n\nDownload: All reports at once as ZIP or individually",
                        slide_colors[5])
    
    # Slide 7: Tech Stack
    add_tech_slide(prs, "🔧 Technology (For IT Staff)",
                  [("Frontend", "Streamlit - Modern web interface, no coding needed"),
                   ("Backend", "Python - Powerful data processing"),
                   ("Data Processing", "Pandas - Excel file reading & manipulation"),
                   ("Document Generation", "Python-docx - Word file creation & customization")],
                  "Deployment: Cloud-based (Streamlit Cloud) - Secure, scalable, automatic updates",
                  slide_colors[6])
    
    # Slide 8: Benefits & Call to Action
    add_benefits_slide(prs, "✅ Benefits",
                      [("⏱️", "Save Time", "Hours of work done in minutes"),
                       ("✓", "Reduce Errors", "Consistent, accurate reports"),
                       ("🔒", "Secure & Private", "Your data stays with you")],
                      "Ready to transform your report card process?",
                      slide_colors[7])
    
    # Save presentation
    output_path = "/Users/christine.anuli/Downloads/report_card_generator/PRESENTATION.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
    return output_path


def add_title_slide(prs, title, subtitle, tagline, colors):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors[0])
    
    # Title
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    left = Inches(0.5)
    top = Inches(3.8)
    width = Inches(9)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(*colors[1])
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    left = Inches(0.5)
    top = Inches(5.2)
    width = Inches(9)
    height = Inches(0.8)
    tagline_box = slide.shapes.add_textbox(left, top, width, height)
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*colors[1])
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullet_points, colors):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors[0])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    
    # Content
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for idx, point in enumerate(bullet_points):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(*colors[1])
        p.level = 0
        p.space_before = Pt(10)


def add_solution_slide(prs, title, description, steps, colors):
    """Add the solution slide with 4 steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors[0])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    
    # Description
    left = Inches(1)
    top = Inches(1.2)
    width = Inches(8)
    height = Inches(0.8)
    desc_box = slide.shapes.add_textbox(left, top, width, height)
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(*colors[1])
    p.alignment = PP_ALIGN.CENTER
    
    # Steps
    step_width = 2
    step_height = 3
    start_left = 1
    top = 2.3
    
    for idx, step in enumerate(steps):
        left = start_left + (idx * (step_width + 0.3))
        # Step box
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), 
                                       Inches(step_width), Inches(step_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.fill.transparency = 0.2
        shape.line.color.rgb = RGBColor(*colors[1])
        
        # Step number
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"{idx + 1}"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors[1])
        p.alignment = PP_ALIGN.CENTER
        
        # Step text
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*colors[1])
        p.alignment = PP_ALIGN.CENTER


def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content, colors):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors[0])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    
    # Left column
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(4.5)
    height = Inches(5.5)
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    
    p = left_frame.add_paragraph()
    p.text = left_content
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*colors[1])
    p.space_before = Pt(10)
    
    # Right column
    left = Inches(5.2)
    top = Inches(1.3)
    width = Inches(4.3)
    height = Inches(5.5)
    right_box = slide.shapes.add_textbox(left, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    
    p = right_frame.add_paragraph()
    p.text = right_content
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(*colors[1])
    p.space_before = Pt(10)


def add_tech_slide(prs, title, tech_items, deployment, colors):
    """Add the technology slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors[0])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    
    # Tech items
    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(8.4)
    height = Inches(4.5)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for idx, (tech, description) in enumerate(tech_items):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = f"{tech}: {description}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*colors[1])
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Deployment info
    p = text_frame.add_paragraph()
    p.text = deployment
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    p.space_before = Pt(15)


def add_benefits_slide(prs, title, benefits, cta, colors):
    """Add the benefits slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*colors[0])
    
    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    
    # Benefits in columns
    benefit_width = 2.8
    benefit_height = 3.5
    start_left = 0.6
    top = 1.4
    
    for idx, (icon, benefit_title, benefit_text) in enumerate(benefits):
        left = start_left + (idx * (benefit_width + 0.3))
        # Benefit box
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                       Inches(benefit_width), Inches(benefit_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.fill.transparency = 0.2
        shape.line.color.rgb = RGBColor(*colors[1])
        
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = icon
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER
        
        p = text_frame.add_paragraph()
        p.text = benefit_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*colors[1])
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(5)
        
        p = text_frame.add_paragraph()
        p.text = benefit_text
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(*colors[1])
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(5)
    
    # CTA
    left = Inches(1)
    top = Inches(5.3)
    width = Inches(8)
    height = Inches(1)
    cta_box = slide.shapes.add_textbox(left, top, width, height)
    cta_frame = cta_box.text_frame
    p = cta_frame.paragraphs[0]
    p.text = cta
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*colors[1])
    p.alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    create_powerpoint()
